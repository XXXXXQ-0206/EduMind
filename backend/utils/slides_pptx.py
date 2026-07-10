"""Build simple PPTX files from generated slide previews."""
from __future__ import annotations

from pathlib import Path
from typing import TYPE_CHECKING, Mapping, Sequence
from urllib.parse import unquote, urlparse

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

if TYPE_CHECKING:
    from infrastructure.object_store import ObjectStore


FONT_NAME = "Microsoft YaHei"


def build_simple_slides_pptx(
    slides: Sequence[Mapping[str, object]],
    output_path: Path,
    *,
    deck_title: str = "",
    object_store: ObjectStore | None = None,
) -> Path:
    """Create a fixed-template PPTX with one image and text block per page."""
    normalized = [_normalize_slide(item, index, deck_title) for index, item in enumerate(slides)]
    if not normalized:
        raise ValueError("slides cannot be empty")

    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for slide_data in normalized:
        page = prs.slides.add_slide(blank_layout)
        _paint_background(page)
        _add_page_number(page, slide_data["page_index"], len(normalized))
        _add_title(page, slide_data["title"])
        _add_image(page, slide_data, object_store)
        _add_bullets(page, slide_data["bullets"], slide_data["title"])

    prs.save(output_path)
    return output_path


def _normalize_slide(item: Mapping[str, object], index: int, deck_title: str) -> dict:
    title = str(item.get("title") or "").strip()
    if not title:
        title = deck_title if index == 0 and deck_title else f"第{index + 1}页"

    raw_bullets = item.get("bullets")
    bullets: list[str] = []
    if isinstance(raw_bullets, list):
        bullets = [str(value).strip() for value in raw_bullets if str(value).strip()]
    elif raw_bullets:
        bullets = [str(raw_bullets).strip()]

    return {
        "page_index": index + 1,
        "title": title,
        "bullets": bullets[:6],
        "imageObjectKey": str(item.get("imageObjectKey") or "").strip(),
        "imageUrl": str(item.get("imageUrl") or "").strip(),
    }


def _paint_background(slide) -> None:
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = RGBColor(248, 250, 252)

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.16))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(37, 99, 235)
    accent.line.fill.background()


def _add_page_number(slide, page_index: int, total: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.65), Inches(6.95), Inches(1.05), Inches(0.25))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = f"{page_index}/{total}"
    paragraph.font.size = Pt(9)
    paragraph.font.name = FONT_NAME
    paragraph.font.color.rgb = RGBColor(100, 116, 139)


def _add_title(slide, title: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.62), Inches(0.42), Inches(12.1), Inches(0.72))
    frame = box.text_frame
    frame.word_wrap = True
    paragraph = frame.paragraphs[0]
    paragraph.text = title[:80]
    paragraph.font.bold = True
    paragraph.font.name = FONT_NAME
    paragraph.font.size = Pt(30 if len(title) <= 24 else 25)
    paragraph.font.color.rgb = RGBColor(15, 23, 42)


def _add_image(slide, slide_data: Mapping[str, object], object_store: ObjectStore | None) -> None:
    left = Inches(0.66)
    top = Inches(1.35)
    width = Inches(6.0)
    height = Inches(5.45)

    frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(226, 232, 240)
    frame.line.color.rgb = RGBColor(203, 213, 225)

    image_path = _resolve_image_path(slide_data, object_store)
    if not image_path or not image_path.is_file():
        _add_image_placeholder(slide, left, top, width, height)
        return

    try:
        from PIL import Image

        with Image.open(image_path) as image:
            image_width, image_height = image.size
        if image_width <= 0 or image_height <= 0:
            raise ValueError("invalid image dimensions")

        scale = min(int(width) / image_width, int(height) / image_height)
        fitted_width = int(image_width * scale)
        fitted_height = int(image_height * scale)
        fitted_left = int(left) + (int(width) - fitted_width) // 2
        fitted_top = int(top) + (int(height) - fitted_height) // 2
        slide.shapes.add_picture(str(image_path), fitted_left, fitted_top, width=fitted_width, height=fitted_height)
    except Exception:
        _add_image_placeholder(slide, left, top, width, height)


def _resolve_image_path(slide_data: Mapping[str, object], object_store: ObjectStore | None) -> Path | None:
    if not object_store:
        return None

    object_key = str(slide_data.get("imageObjectKey") or "").strip()
    if object_key:
        return object_store.path_for(object_key)

    image_url = str(slide_data.get("imageUrl") or "").strip()
    object_key = _object_key_from_url(image_url)
    if object_key:
        return object_store.path_for(object_key)

    if image_url:
        candidate = Path(image_url)
        if candidate.is_file():
            return candidate

    return None


def _object_key_from_url(value: str) -> str:
    if not value:
        return ""

    parsed = urlparse(value)
    path = unquote(parsed.path if parsed.scheme else value)
    public_base = "/storage"
    try:
        from config import config

        public_base = (config.object_store_base_url or public_base).rstrip("/")
    except Exception:
        pass
    if path.startswith(f"{public_base}/"):
        return path[len(public_base) + 1 :]
    return ""


def _add_image_placeholder(slide, left, top, width, height) -> None:
    box = slide.shapes.add_textbox(left, top + height // 2 - Inches(0.18), width, Inches(0.4))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = "配图暂未生成"
    paragraph.font.name = FONT_NAME
    paragraph.font.size = Pt(18)
    paragraph.font.color.rgb = RGBColor(100, 116, 139)


def _add_bullets(slide, bullets: Sequence[str], title: str) -> None:
    box = slide.shapes.add_textbox(Inches(7.12), Inches(1.42), Inches(5.55), Inches(5.25))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.05)
    frame.margin_bottom = Inches(0.05)
    frame.clear()

    content = list(bullets) or [title]
    font_size = 21 if len(content) <= 3 else 18 if len(content) <= 5 else 16

    for index, bullet in enumerate(content):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = bullet[:160]
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        paragraph.font.name = FONT_NAME
        paragraph.font.size = Pt(font_size)
        paragraph.font.color.rgb = RGBColor(30, 41, 59)
