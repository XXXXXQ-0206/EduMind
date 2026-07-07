"""
文档解析工具
支持 PDF, DOCX, PPTX, XLSX, Markdown, TXT 与常见纯文本文件解析。
"""
import asyncio
import aiofiles
from pathlib import Path
from typing import Optional
import re
import shutil
import tempfile
import uuid
import pymupdf  # PyMuPDF
import markdownify
import mammoth

from config import config


def _storage_root() -> Path:
    return config.storage_dir.resolve()


def _ensure_storage_file(file_path: str | Path) -> Path:
    path = Path(file_path)
    resolved = path.resolve()
    root = _storage_root()
    if resolved != root and root not in resolved.parents:
        raise ValueError("File path is outside the configured storage directory")
    if not resolved.exists():
        raise FileNotFoundError(f"File not found: {path.name}")
    if not resolved.is_file():
        raise ValueError("File path does not point to a regular file")
    return resolved


def _safe_upload_filename(filename: str | None) -> str:
    name = Path(filename or "upload.bin").name.strip() or "upload.bin"
    cleaned = re.sub(r"[^A-Za-z0-9._\-\u4e00-\u9fff]+", "_", name).strip("._")
    return cleaned[:180] or "upload.bin"


async def extract_text_from_file(
    file_path: str,
    mime_type: Optional[str] = None,
) -> str:
    """
    从文件中提取文本

    Args:
        file_path: 文件路径
        mime_type: MIME 类型（可选）

    Returns:
        提取的文本内容
    """
    path = _ensure_storage_file(file_path)

    # 如果没有提供 MIME 类型，根据文件扩展名推断
    if not mime_type:
        mime_type = _guess_mime_type(path)

    # PDF
    if "pdf" in mime_type.lower():
        return await _extract_pdf(path)

    lower_mime = mime_type.lower()
    suffix = path.suffix.lower()

    # DOCX
    if "wordprocessingml" in lower_mime or suffix == ".docx":
        return await _extract_docx(path)

    # 旧版 DOC：优先用 Word/LibreOffice 转换后提取
    if "msword" in lower_mime or suffix == ".doc":
        return await _extract_legacy_doc(path)

    # PPTX
    if "presentationml" in lower_mime or suffix == ".pptx":
        return await _extract_pptx(path)

    # XLSX
    if "spreadsheetml" in lower_mime or suffix == ".xlsx":
        return await _extract_xlsx(path)

    # Markdown
    if "markdown" in lower_mime or suffix in [".md", ".markdown"]:
        return await _extract_markdown(path)

    # HTML
    if "html" in lower_mime or suffix in [".html", ".htm"]:
        return await _extract_html(path)

    # TXT
    if "text" in lower_mime or _looks_like_text_file(path):
        return await _extract_text(path)

    # 最后尝试按文本文件解码，便于处理 .csv/.json/.py/.ts 等未声明 MIME 的文本资料。
    try:
        return await _extract_text(path)
    except UnicodeDecodeError as exc:
        raise ValueError(f"Unsupported file type: {mime_type}") from exc


def _guess_mime_type(path: Path) -> str:
    """根据文件扩展名推断 MIME 类型"""
    ext = path.suffix.lower()

    mime_map = {
        ".pdf": "application/pdf",
        ".doc": "application/msword",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".markdown": "text/markdown",
        ".csv": "text/csv",
        ".tsv": "text/tab-separated-values",
        ".json": "application/json",
        ".yaml": "text/yaml",
        ".yml": "text/yaml",
        ".html": "text/html",
        ".htm": "text/html",
    }

    return mime_map.get(ext, "application/octet-stream")


async def _extract_pdf(file_path: str | Path) -> str:
    """从 PDF 提取文本"""
    path = _ensure_storage_file(file_path)
    doc = pymupdf.open(str(path))
    text_parts = []

    for page in doc:
        text_parts.append(page.get_text())

    doc.close()
    return "\n\n".join(text_parts)


async def _extract_docx(file_path: str | Path) -> str:
    """从 DOCX 提取文本"""
    path = _ensure_storage_file(file_path)
    # 使用 mammoth 提取原始文本
    with open(path, "rb") as docx_file:
        result = mammoth.extract_raw_text(docx_file)
        return result.value


async def _extract_legacy_doc(file_path: str | Path) -> str:
    """从旧版 .doc 提取文本。"""
    path = _ensure_storage_file(file_path)
    try:
        return await _convert_office_document(path, target_suffix=".docx", extractor=_extract_docx)
    except Exception as exc:
        raise ValueError(f"Unsupported file type: {path.suffix}") from exc


async def _extract_pptx(file_path: str | Path) -> str:
    """从 PPTX 提取文本。"""
    from pptx import Presentation

    path = _ensure_storage_file(file_path)
    prs = Presentation(str(path))
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text:
                parts.append(text.strip())
        if parts:
            slides.append(f"第 {index} 页\n" + "\n".join(parts))
    return "\n\n".join(slides)


async def _extract_xlsx(file_path: str | Path) -> str:
    """从 XLSX 提取文本。"""
    from openpyxl import load_workbook

    path = _ensure_storage_file(file_path)
    workbook = load_workbook(path, data_only=True, read_only=True)
    sheets = []
    for sheet in workbook.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets.append(f"工作表 {sheet.title}\n" + "\n".join(rows))
    workbook.close()
    return "\n\n".join(sheets)


async def _extract_markdown(file_path: str | Path) -> str:
    """从 Markdown 提取文本（转换为纯文本）"""
    path = _ensure_storage_file(file_path)
    async with aiofiles.open(path, "r", encoding="utf-8") as f:
        content = await f.read()

    # 使用 markdownify 将 markdown 转换为纯文本
    # 或者直接返回 markdown 内容
    return content


async def _extract_html(file_path: str | Path) -> str:
    """从 HTML 提取可读文本。"""
    path = _ensure_storage_file(file_path)
    async with aiofiles.open(path, "r", encoding="utf-8") as f:
        content = await f.read()
    return markdownify.markdownify(content, heading_style="ATX")


async def _extract_text(file_path: str | Path) -> str:
    """从 TXT 提取文本"""
    path = _ensure_storage_file(file_path)
    for encoding in ("utf-8", "utf-8-sig", "gb18030", "gbk", "latin-1"):
        try:
            async with aiofiles.open(path, "r", encoding=encoding) as f:
                return await f.read()
        except UnicodeDecodeError:
            continue
    async with aiofiles.open(path, "rb") as f:
        return (await f.read()).decode("utf-8", errors="ignore")


def _looks_like_text_file(path: Path) -> bool:
    """判断未知 MIME 的文件是否属于常见可文本化资料。"""
    return path.suffix.lower() in {
        ".txt",
        ".csv",
        ".tsv",
        ".json",
        ".jsonl",
        ".yaml",
        ".yml",
        ".xml",
        ".html",
        ".htm",
        ".css",
        ".js",
        ".ts",
        ".vue",
        ".py",
        ".java",
        ".c",
        ".cpp",
        ".h",
        ".hpp",
        ".md",
        ".markdown",
        ".log",
    }


async def _convert_office_document(path: Path, *, target_suffix: str, extractor):
    """尝试把旧 Office 文件转成可抽取格式。"""
    source = _ensure_storage_file(path)
    if source.suffix.lower() == target_suffix:
        return await extractor(source)

    if target_suffix not in {".docx"}:
        raise ValueError("Unsupported conversion target")

    temp_parent = (_storage_root() / ".tmp" / "office-convert").resolve()
    temp_parent.mkdir(parents=True, exist_ok=True)
    with tempfile.TemporaryDirectory(prefix="edumind-office-", dir=str(temp_parent)) as temp_name:
        temp_dir = Path(temp_name).resolve()
        temp_source = temp_dir / f"input{source.suffix.lower()}"
        converted = temp_dir / f"input{target_suffix}"
        shutil.copyfile(source, temp_source)

        # 1) 先尝试调用 Word COM
        try:
            import win32com.client  # type: ignore

            try:
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(str(temp_source), False, True)
                try:
                    doc.SaveAs2(str(converted), FileFormat=12)
                finally:
                    doc.Close(False)
                    word.Quit()
                if converted.exists():
                    return await extractor(converted)
            finally:
                if converted.exists():
                    converted.unlink()
        except Exception:
            pass

        # 2) 再尝试 LibreOffice/soffice
        try:
            import subprocess

            soffice = shutil.which("soffice") or shutil.which("libreoffice")
            if not soffice:
                raise FileNotFoundError("LibreOffice executable not found")
            cmd = [
                soffice,
                "--headless",
                "--convert-to",
                target_suffix.lstrip("."),
                "--outdir",
                str(temp_dir),
                "--",
                str(temp_source),
            ]
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            if converted.exists():
                return await extractor(converted)
        except Exception:
            pass

    raise FileNotFoundError(f"Could not convert legacy office file: {source.name}")


async def parse_multipart_form(
    files: list,
    form_data: dict,
) -> dict:
    """
    解析 multipart 表单数据

    Args:
        files: 文件列表
        form_data: 表单数据字典

    Returns:
        解析后的数据字典
    """
    result = {
        "q": form_data.get("q", ""),
        "chatId": form_data.get("chatId"),
        "files": [],
    }

    upload_dir = (_storage_root() / "uploads").resolve()
    upload_dir.mkdir(parents=True, exist_ok=True)

    for file in files:
        filename = _safe_upload_filename(file.filename)
        content_type = file.content_type

        # 保存文件
        file_path = (upload_dir / f"{uuid.uuid4().hex}-{filename}").resolve()
        if upload_dir not in file_path.parents:
            raise ValueError("Invalid upload filename")

        async with aiofiles.open(file_path, "wb") as f:
            content = await file.read()
            await f.write(content)

        result["files"].append({
            "path": str(file_path),
            "filename": filename,
            "mimeType": content_type,
        })

    return result
